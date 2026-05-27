"""
Build the LLMai hackathon submission deck.

Run:    python scripts/build_hackathon_deck.py
Output: docs/LLMai-hackathon-deck.pptx  (11 slides, 16:9 widescreen, dark theme)
"""
from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt, Emu

# ── Palette ──────────────────────────────────────────────────────────────────
BG       = RGBColor(0x07, 0x09, 0x0C)
TEXT     = RGBColor(0xF5, 0xF7, 0xFA)
GREEN    = RGBColor(0x46, 0xD3, 0x9A)
BLUE     = RGBColor(0x6A, 0xA9, 0xFF)
MUTED    = RGBColor(0xAA, 0xB4, 0xC0)
QUIET    = RGBColor(0x75, 0x82, 0x93)
LINE     = RGBColor(0x26, 0x31, 0x3D)
CARD_BG  = RGBColor(0x10, 0x14, 0x19)
CODE_BG  = RGBColor(0x0A, 0x0D, 0x11)
AMBER    = RGBColor(0xF6, 0xC5, 0x6F)

# Typography
SANS = "Inter"          # falls back to Calibri on systems without Inter
MONO = "Consolas"       # universally available on Windows; close enough on mac/linux

# Layout (13.333 × 7.5)
W = 13.333
H = 7.5
MARGIN = 0.6

# ── Presentation setup ──────────────────────────────────────────────────────

prs = Presentation()
prs.slide_width = Inches(W)
prs.slide_height = Inches(H)
BLANK = prs.slide_layouts[6]   # truly blank


# ── Helpers ──────────────────────────────────────────────────────────────────

def add_slide():
    s = prs.slides.add_slide(BLANK)
    bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.line.fill.background()
    bg.fill.solid()
    bg.fill.fore_color.rgb = BG
    bg.shadow.inherit = False
    return s


def add_text(slide, text, *, x, y, w, h, size=18, color=TEXT, bold=False,
             font=SANS, align=PP_ALIGN.LEFT, valign=MSO_ANCHOR.TOP, italic=False):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = 0
    tf.margin_top = tf.margin_bottom = 0
    tf.vertical_anchor = valign
    # First paragraph
    if isinstance(text, str):
        runs = [(text, color, size, bold, font, italic)]
    else:
        runs = text  # list of tuples
    for i, run in enumerate(runs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        if isinstance(run, tuple):
            txt, c, sz, b, fnt, it = (list(run) + [italic])[:6]
        else:
            txt, c, sz, b, fnt, it = run, color, size, bold, font, italic
        r = p.add_run()
        r.text = txt
        r.font.name = fnt
        r.font.size = Pt(sz)
        r.font.bold = b
        r.font.italic = it
        r.font.color.rgb = c
    return box


def add_paragraphs(slide, lines, *, x, y, w, h, size=14, color=TEXT, font=SANS,
                   line_spacing=1.25, bullet=False, valign=MSO_ANCHOR.TOP):
    """Add a multi-line text box. Each line can be a string or (text, options-dict)."""
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = 0
    tf.margin_top = tf.margin_bottom = 0
    tf.vertical_anchor = valign
    for i, line in enumerate(lines):
        if isinstance(line, str):
            txt, opts = line, {}
        else:
            txt, opts = line
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = opts.get("align", PP_ALIGN.LEFT)
        p.line_spacing = line_spacing
        if i > 0:
            p.space_before = Pt(opts.get("space_before", 6))
        text_with_bullet = (("• " if bullet and not opts.get("no_bullet") else "") + txt)
        r = p.add_run()
        r.text = text_with_bullet
        r.font.name = opts.get("font", font)
        r.font.size = Pt(opts.get("size", size))
        r.font.bold = opts.get("bold", False)
        r.font.italic = opts.get("italic", False)
        r.font.color.rgb = opts.get("color", color)
    return box


def add_rect(slide, x, y, w, h, *, fill=CARD_BG, line=None, line_width=0.75):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.shadow.inherit = False
    if line is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = line
        shape.line.width = Pt(line_width)
    return shape


def add_accent_bar(slide, x, y, w, h, *, color=GREEN):
    """Solid colored bar (used for top/side accents on cards)."""
    return add_rect(slide, x, y, w, h, fill=color)


def add_section_title(slide, text, *, y=0.55):
    """Standard slide title in accent green."""
    add_text(slide, text, x=MARGIN, y=y, w=W - 2 * MARGIN, h=0.7,
             size=28, bold=True, color=GREEN, font=SANS)


def add_page_footer(slide, n, total):
    add_text(slide, f"{n} / {total}", x=W - 1.5, y=H - 0.5, w=1.0, h=0.3,
             size=10, color=QUIET, font=MONO, align=PP_ALIGN.RIGHT)
    add_text(slide, "github.com/sechan9999/LLMai", x=MARGIN, y=H - 0.5,
             w=6, h=0.3, size=10, color=QUIET, font=MONO)


# ── Slides ───────────────────────────────────────────────────────────────────

TOTAL = 11


def slide_1_title():
    s = add_slide()
    # Brand mark — small green/blue gradient hint via two stacked rects
    add_rect(s, MARGIN, MARGIN, 0.5, 0.5, fill=GREEN)
    add_text(s, "L", x=MARGIN, y=MARGIN - 0.02, w=0.5, h=0.5,
             size=24, bold=True, color=BG, font=SANS, align=PP_ALIGN.CENTER,
             valign=MSO_ANCHOR.MIDDLE)

    # Title block
    add_text(s, "LLMai", x=MARGIN, y=2.35, w=W - 2 * MARGIN, h=1.85,
             size=92, bold=True, color=TEXT, font=SANS, align=PP_ALIGN.LEFT)
    add_text(s, "A local-first AI coding agent with three layers of awareness",
             x=MARGIN, y=4.25, w=W - 2 * MARGIN, h=0.7,
             size=22, color=MUTED, font=SANS)

    # Underline accent (just under the subtitle)
    add_rect(s, MARGIN, 5.05, 1.6, 0.04, fill=GREEN)

    # Footer
    add_text(s, "Hackathon submission · 2026", x=MARGIN, y=H - 0.6, w=5, h=0.3,
             size=12, color=QUIET, font=MONO)
    add_text(s, "github.com/sechan9999/LLMai", x=W - 5 - MARGIN, y=H - 0.6,
             w=5, h=0.3, size=12, color=QUIET, font=MONO, align=PP_ALIGN.RIGHT)


def slide_2_problem():
    s = add_slide()
    add_section_title(s, "Every AI coding agent has the same two blind spots")

    col_w = (W - 2 * MARGIN - 0.5) / 2
    y_top = 1.9
    h_card = 4.2

    for i, (heading, body, accent) in enumerate([
        ("Forgets yesterday",
         "Every conversation starts from zero. Your past decisions, debug sessions, "
         "and the rationale behind your architecture vanish at /reset.",
         GREEN),
        ("Ignores what your team already learned",
         "The model writes a fix while the issue your team filed six months ago — "
         "with the exact same fix — sits unread in GitLab.",
         BLUE),
    ]):
        x = MARGIN + i * (col_w + 0.5)
        add_rect(s, x, y_top, col_w, h_card, fill=CARD_BG)
        add_accent_bar(s, x, y_top, col_w, 0.08, color=accent)
        add_text(s, heading, x=x + 0.4, y=y_top + 0.4, w=col_w - 0.8, h=0.7,
                 size=22, bold=True, color=TEXT)
        add_text(s, body, x=x + 0.4, y=y_top + 1.3, w=col_w - 0.8, h=h_card - 1.5,
                 size=15, color=MUTED)

    add_page_footer(s, 2, TOTAL)


def slide_3_solution():
    s = add_slide()
    add_section_title(s, "LLMai gives the agent three distinct kinds of awareness")

    col_w = (W - 2 * MARGIN - 1.0) / 3
    y_top = 1.9
    h_card = 4.0

    cols = [
        ("Operational", "Dynatrace",
         "Every tool call traced. Latency, tokens, permission outcome, success / error."),
        ("Personal", "MongoDB Atlas",
         "Remembers your past sessions per workspace. Boots warm with prior summaries auto-injected."),
        ("Organizational", "Elastic",
         "Hybrid search over GitLab issues, CI logs, docs. Calls search_knowledge before writing code."),
    ]
    for i, (kind, partner, desc) in enumerate(cols):
        x = MARGIN + i * (col_w + 0.5)
        add_rect(s, x, y_top, col_w, h_card, fill=CARD_BG)
        add_accent_bar(s, x, y_top, col_w, 0.08, color=GREEN)
        # Kind label in muted mono
        add_text(s, kind.upper(), x=x + 0.4, y=y_top + 0.35, w=col_w - 0.8, h=0.3,
                 size=11, bold=True, color=QUIET, font=MONO)
        # Partner name in accent blue
        add_text(s, partner, x=x + 0.4, y=y_top + 0.7, w=col_w - 0.8, h=0.7,
                 size=24, bold=True, color=BLUE)
        # Description
        add_text(s, desc, x=x + 0.4, y=y_top + 1.7, w=col_w - 0.8, h=h_card - 1.9,
                 size=14, color=MUTED)

    # Bottom strip
    add_text(s, "All three are opt-in. Core agent runs 100% locally.",
             x=MARGIN, y=H - 1.1, w=W - 2 * MARGIN, h=0.4,
             size=14, color=GREEN, font=MONO, align=PP_ALIGN.CENTER, italic=True)

    add_page_footer(s, 3, TOTAL)


def slide_4_architecture():
    s = add_slide()
    add_section_title(s, "How the three layers wire into the agent loop")

    diagram = (
        "                                       search_knowledge   ──►  Elastic  (issues + docs)\n"
        "                                       │\n"
        "    Browser / CLI ──► Agent Loop ──────┼─► recall_memory     ──►  MongoDB Atlas\n"
        "                          │            │                          (sessions + summaries\n"
        "                          │            └─► query_logs        ──►   + extracted knowledge)\n"
        "                          │\n"
        "                          ▼\n"
        "                     OTel spans\n"
        "                          │\n"
        "                          ▼\n"
        "                  Bindplane ──┬──► Dynatrace             (traces + metrics)\n"
        "                              │\n"
        "                              └──► Elastic agent-logs    (agent queries itself)\n"
        "\n"
        "    LLM backend:   Ollama localhost:11434   (Gemini / Groq optional)"
    )
    # Code-like block
    add_rect(s, MARGIN, 1.6, W - 2 * MARGIN, 4.7, fill=CODE_BG, line=LINE)
    add_text(s, diagram, x=MARGIN + 0.3, y=1.75, w=W - 2 * MARGIN - 0.6, h=4.5,
             size=12, color=TEXT, font=MONO)

    add_text(s, "Each layer is independent. Failure in one cannot cascade to another.",
             x=MARGIN, y=6.5, w=W - 2 * MARGIN, h=0.4,
             size=13, color=MUTED, italic=True, align=PP_ALIGN.CENTER)

    add_page_footer(s, 4, TOTAL)


def slide_5_dynatrace():
    s = add_slide()
    add_section_title(s, "Layer 1 — Operational awareness via OpenTelemetry")
    add_text(s, "DYNATRACE", x=MARGIN, y=1.3, w=4, h=0.3, size=11, bold=True,
             color=BLUE, font=MONO)

    col_w = (W - 2 * MARGIN - 0.5) / 2
    y_top = 1.95
    h_card = 4.4

    # Left card — What's instrumented
    add_rect(s, MARGIN, y_top, col_w, h_card, fill=CARD_BG)
    add_accent_bar(s, MARGIN, y_top, col_w, 0.06, color=GREEN)
    add_text(s, "What's instrumented", x=MARGIN + 0.4, y=y_top + 0.3, w=col_w - 0.8,
             h=0.5, size=18, bold=True, color=TEXT)
    spans = [
        ("agent.turn", "one per user input"),
        ("agent.iteration", "per loop iteration  (≤ 20)"),
        ("llm.chat", "per LLM call  — model, tokens, latency"),
        ("tool.invocation", "name, args preview, permission outcome, exec latency, result chars"),
    ]
    y_cursor = y_top + 1.05
    for name, desc in spans:
        add_text(s, name, x=MARGIN + 0.4, y=y_cursor, w=col_w - 0.8, h=0.3,
                 size=13, color=GREEN, font=MONO, bold=True)
        add_text(s, desc, x=MARGIN + 0.4, y=y_cursor + 0.3, w=col_w - 0.8, h=0.5,
                 size=12, color=MUTED)
        y_cursor += 0.8

    # Right card — Privacy
    x_right = MARGIN + col_w + 0.5
    add_rect(s, x_right, y_top, col_w, h_card, fill=CARD_BG)
    add_accent_bar(s, x_right, y_top, col_w, 0.06, color=GREEN)
    add_text(s, "Privacy", x=x_right + 0.4, y=y_top + 0.3, w=col_w - 0.8, h=0.5,
             size=18, bold=True, color=TEXT)
    privacy = [
        "No raw prompts in spans",
        "No file contents in attributes",
        "tool.args.preview truncated to 200 chars",
        "Opt-in via LLMAI_OTEL_ENABLED=true",
    ]
    add_paragraphs(s, privacy, x=x_right + 0.4, y=y_top + 1.1, w=col_w - 0.8,
                   h=h_card - 1.3, size=14, color=MUTED, bullet=True, line_spacing=1.4)

    add_page_footer(s, 5, TOTAL)


def slide_6_atlas():
    s = add_slide()
    add_section_title(s, "Layer 2 — Cross-session memory via Atlas Vector Search")
    add_text(s, "MONGODB ATLAS", x=MARGIN, y=1.3, w=4, h=0.3, size=11, bold=True,
             color=BLUE, font=MONO)

    rows = [
        ("sessions",  "Raw transcripts per turn"),
        ("summaries", "LLM-summarized, vector-embedded  (cosine, 768d)  for cross-session recall"),
        ("knowledge", "Extracted facts  (3–5 per session), vector-embedded, semantically searchable"),
    ]
    y_top = 1.95
    row_h = 0.85
    for i, (name, desc) in enumerate(rows):
        y = y_top + i * (row_h + 0.18)
        add_rect(s, MARGIN, y, W - 2 * MARGIN, row_h, fill=CARD_BG)
        add_accent_bar(s, MARGIN, y, 0.08, row_h, color=GREEN)
        add_text(s, name, x=MARGIN + 0.4, y=y + 0.18, w=2.6, h=0.5,
                 size=18, bold=True, color=GREEN, font=MONO)
        add_text(s, desc, x=MARGIN + 3.2, y=y + 0.22, w=W - 2 * MARGIN - 3.6, h=0.5,
                 size=14, color=MUTED, valign=MSO_ANCHOR.TOP)

    # Scoping note
    add_text(s, "Per-workspace scoping via sha256(absolute_path)[:16] — survives directory renames.",
             x=MARGIN, y=5.05, w=W - 2 * MARGIN, h=0.4, size=13, color=MUTED, italic=True)

    # Code block — recall_memory example
    code = ('recall_memory(query="rate limiting bug in /api/chat",\n'
            '              scope="both", limit=5)')
    add_rect(s, MARGIN, 5.6, W - 2 * MARGIN, 1.1, fill=CODE_BG, line=LINE)
    add_text(s, code, x=MARGIN + 0.3, y=5.7, w=W - 2 * MARGIN - 0.6, h=0.9,
             size=14, color=GREEN, font=MONO)

    add_page_footer(s, 6, TOTAL)


def slide_7_elastic():
    s = add_slide()
    add_section_title(s, "Layer 3 — Hybrid search over org knowledge via Elasticsearch")
    add_text(s, "ELASTIC", x=MARGIN, y=1.3, w=4, h=0.3, size=11, bold=True,
             color=BLUE, font=MONO)

    col_w = (W - 2 * MARGIN - 0.5) / 2
    y_top = 1.95
    h_card = 4.4

    # Left — Two tools exposed
    add_rect(s, MARGIN, y_top, col_w, h_card, fill=CARD_BG)
    add_accent_bar(s, MARGIN, y_top, col_w, 0.06, color=GREEN)
    add_text(s, "Two tools exposed", x=MARGIN + 0.4, y=y_top + 0.3, w=col_w - 0.8,
             h=0.5, size=18, bold=True, color=TEXT)
    tools = [
        ("search_knowledge",
         "Hybrid BM25 + kNN, RRF on Atlas / Cloud, kNN-only fallback on basic license. "
         "Auto-approved."),
        ("query_logs",
         "Raw ES|QL over pipeline failures + the agent's own self-logs. Permission-gated."),
    ]
    y_cursor = y_top + 1.05
    for name, desc in tools:
        add_text(s, name, x=MARGIN + 0.4, y=y_cursor, w=col_w - 0.8, h=0.35,
                 size=14, color=GREEN, font=MONO, bold=True)
        add_text(s, desc, x=MARGIN + 0.4, y=y_cursor + 0.4, w=col_w - 0.8, h=1.4,
                 size=13, color=MUTED)
        y_cursor += 1.75

    # Right — Why this matters
    x_right = MARGIN + col_w + 0.5
    add_rect(s, x_right, y_top, col_w, h_card, fill=CARD_BG)
    add_accent_bar(s, x_right, y_top, col_w, 0.06, color=GREEN)
    add_text(s, "Why this matters", x=x_right + 0.4, y=y_top + 0.3, w=col_w - 0.8,
             h=0.5, size=18, bold=True, color=TEXT)
    why = [
        "Agent calls search_knowledge BEFORE writing code that touches an error path",
        "Surfaces the GitLab issue from 6 months ago instead of inventing a fix",
        "Via Bindplane tee, agent can query its own behavior with ES|QL",
    ]
    add_paragraphs(s, why, x=x_right + 0.4, y=y_top + 1.1, w=col_w - 0.8,
                   h=h_card - 1.3, size=14, color=MUTED, bullet=True, line_spacing=1.45)

    add_page_footer(s, 7, TOTAL)


def slide_8_verified():
    s = add_slide()
    add_section_title(s, "Verified end-to-end")
    add_text(s, "What works today, tested locally", x=MARGIN, y=1.3, w=10, h=0.3,
             size=13, color=MUTED, italic=True)

    rows = [
        ('search_knowledge("chat endpoint throttling")',
         'score 0.84 on rate-limit issue  (pure semantic match, no keyword overlap)'),
        ('search_knowledge("cookie token rotation")',
         'score 0.87 on auth design doc'),
        ('query_logs(ES|QL aggregation)',
         'correct row counts grouped by error_signature'),
        ('OTel console exporter',
         'agent.turn → agent.iteration → llm.chat + tool.invocation spans emitted'),
        ('103 unit tests',
         'all pass — graceful degradation verified for every failure mode'),
    ]
    y_top = 1.95
    row_h = 0.78
    for i, (call, result) in enumerate(rows):
        y = y_top + i * (row_h + 0.1)
        add_rect(s, MARGIN, y, W - 2 * MARGIN, row_h, fill=CARD_BG)
        add_accent_bar(s, MARGIN, y, 0.08, row_h, color=GREEN)
        add_text(s, call, x=MARGIN + 0.4, y=y + 0.12, w=6.2, h=0.55,
                 size=13, color=GREEN, font=MONO, valign=MSO_ANCHOR.MIDDLE)
        add_text(s, result, x=MARGIN + 6.8, y=y + 0.12, w=W - MARGIN - 6.8 - 0.4, h=0.55,
                 size=13, color=TEXT, valign=MSO_ANCHOR.MIDDLE)

    add_page_footer(s, 8, TOTAL)


def slide_9_stack():
    s = add_slide()
    add_section_title(s, "Tech stack")

    col_w = (W - 2 * MARGIN - 0.5) / 2
    y_top = 1.7
    h_card = 4.6

    # Left — Core
    add_rect(s, MARGIN, y_top, col_w, h_card, fill=CARD_BG)
    add_accent_bar(s, MARGIN, y_top, col_w, 0.06, color=GREEN)
    add_text(s, "Core", x=MARGIN + 0.4, y=y_top + 0.3, w=col_w - 0.8, h=0.5,
             size=20, bold=True, color=TEXT)
    core = [
        "Python · FastAPI · WebSocket streaming",
        "Ollama  —  qwen2.5-coder, nomic-embed-text",
        "Agent loop:  observe → judge → act   (≤ 20 iterations)",
        "Permission gates:  allow / ask / deny  per tool",
    ]
    add_paragraphs(s, core, x=MARGIN + 0.4, y=y_top + 1.1, w=col_w - 0.8,
                   h=h_card - 1.3, size=14, color=MUTED, bullet=True, line_spacing=1.55)

    # Right — Partner backends
    x_right = MARGIN + col_w + 0.5
    add_rect(s, x_right, y_top, col_w, h_card, fill=CARD_BG)
    add_accent_bar(s, x_right, y_top, col_w, 0.06, color=BLUE)
    add_text(s, "Partner backends  (all opt-in)", x=x_right + 0.4, y=y_top + 0.3,
             w=col_w - 0.8, h=0.5, size=20, bold=True, color=TEXT)
    partners = [
        "Dynatrace  via OpenTelemetry + Bindplane OTel collector",
        "MongoDB Atlas + Vector Search   (768d cosine)",
        "Elasticsearch 8.x  —  RRF, kNN, ES|QL",
        "MCP-compatible tool shapes throughout",
    ]
    add_paragraphs(s, partners, x=x_right + 0.4, y=y_top + 1.1, w=col_w - 0.8,
                   h=h_card - 1.3, size=14, color=MUTED, bullet=True, line_spacing=1.55)

    add_page_footer(s, 9, TOTAL)


def slide_10_roadmap():
    s = add_slide()
    add_section_title(s, "Roadmap")

    items = [
        "Migrate recall_memory + search_knowledge to real MCP transport  (one-file change)",
        "Continuous ingest:  Elastic Agent / Logstash for GitLab + log streams",
        "Cross-workspace recall mode with explicit user opt-in",
        "Auto-route compute:  cheap models for tool selection, larger models for code generation",
        "Add Claude / GPT cloud paths alongside the existing Gemini / Groq options",
    ]
    y_top = 1.9
    item_h = 0.8
    for i, txt in enumerate(items):
        y = y_top + i * (item_h + 0.12)
        # Number circle
        add_rect(s, MARGIN, y, 0.5, 0.5, fill=GREEN)
        add_text(s, str(i + 1), x=MARGIN, y=y, w=0.5, h=0.5, size=18, bold=True,
                 color=BG, font=SANS, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
        # Item text
        add_text(s, txt, x=MARGIN + 0.85, y=y + 0.05, w=W - MARGIN - 0.85 - MARGIN, h=0.5,
                 size=15, color=TEXT, valign=MSO_ANCHOR.MIDDLE)

    add_page_footer(s, 10, TOTAL)


def slide_11_closing():
    s = add_slide()
    # Big tagline
    add_text(s, "Three layers · One agent · Zero API keys required",
             x=MARGIN, y=2.3, w=W - 2 * MARGIN, h=1.6, size=42, bold=True,
             color=TEXT, font=SANS, align=PP_ALIGN.CENTER)
    # Green accent under
    add_rect(s, W / 2 - 0.8, 3.95, 1.6, 0.04, fill=GREEN)
    # Links
    add_text(s, "github.com/sechan9999/LLMai", x=MARGIN, y=4.3, w=W - 2 * MARGIN, h=0.5,
             size=22, color=GREEN, font=MONO, align=PP_ALIGN.CENTER)
    add_text(s, "Live demo:  ll-mai.vercel.app", x=MARGIN, y=4.9, w=W - 2 * MARGIN,
             h=0.5, size=18, color=BLUE, font=MONO, align=PP_ALIGN.CENTER)
    # Run command
    add_text(s, "make demo-up && llmai-server", x=MARGIN, y=5.95, w=W - 2 * MARGIN,
             h=0.5, size=14, color=MUTED, font=MONO, align=PP_ALIGN.CENTER, italic=True)


# ── Build ────────────────────────────────────────────────────────────────────

def main():
    slide_1_title()
    slide_2_problem()
    slide_3_solution()
    slide_4_architecture()
    slide_5_dynatrace()
    slide_6_atlas()
    slide_7_elastic()
    slide_8_verified()
    slide_9_stack()
    slide_10_roadmap()
    slide_11_closing()

    out = Path(__file__).resolve().parent.parent / "docs" / "LLMai-hackathon-deck.pptx"
    out.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(out))
    print(f"wrote {out}  ({len(prs.slides)} slides)")


if __name__ == "__main__":
    main()
